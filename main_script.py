import json
import os
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
import re

# --- CONFIGURATION ---
load_dotenv()


LLM_PROVIDER = os.getenv("LLM_PROVIDER", "GEMINI")
SEARCH_PROVIDER = os.getenv("SEARCH_PROVIDER", "TAVILY") 

# --- 1. TOPIC INPUT ---
def get_topic_from_user():
    """Prompts the user to enter a topic for the presentation."""
    while True:
        topic = input("Enter the topic for your presentation (e.g., 'Climate Change', 'LLM Evaluation'): ").strip()
        if topic:
            print(f"\nTopic selected: \"{topic}\"")
            return topic
        else:
            print("Topic cannot be empty. Please enter a valid topic.")

# --- 2. WEB SEARCH ---
def perform_web_search(topic, search_provider=SEARCH_PROVIDER):
    """
    Performs a web search using the configured search provider.
    This function should return a list of relevant text snippets or summaries.
    """
    print(f"\n[WEB SEARCH - Using {search_provider}] Performing web search for: \"{topic}\"...")
    
    if search_provider == "TAVILY":
        tavily_api_key = os.getenv("TAVILY_API_KEY")
        if tavily_api_key:
            try:
                from tavily import TavilyClient 
                tavily = TavilyClient(api_key=tavily_api_key)
                print(f"[WEB SEARCH - {search_provider}] Sending search query to Tavily API...")
                response = tavily.search(query=f"Key information and recent developments on {topic}", search_depth="advanced", max_results=20)
                snippets = [result['content'] for result in response.get('results', [])]
                print(f"[WEB SEARCH - {search_provider}] Found {len(snippets)} snippets.")
                if not snippets:
                    print(f"[WEB SEARCH - {search_provider}] No snippets found. Returning empty list.")
                    return []
                return snippets
            except ImportError:
                print(f"[WEB SEARCH - {search_provider} Error] The 'tavily-python' library is not installed. Please install it using 'pip install tavily-python'.")
            except Exception as e:
                print(f"[WEB SEARCH - {search_provider} Error] Could not fetch search results: {e}")
                
        else:
            print(f"[WEB SEARCH - {search_provider} Error] TAVILY_API_KEY not found in .env file. Falling back to mock results.")
    
    
    print(f"[WEB SEARCH - Fallback] Using mock search results.")
    mock_results = [
        f"Recent study on '{topic}' shows increasing trends in area X.",
        f"'{topic}' is impacting global markets significantly, especially sector Y.",
        f"Key challenges in '{topic}' include A, B, and C, according to expert Z.",
        f"Innovations in '{topic}' are driven by new technologies like AI and blockchain.",
        f"Future outlook for '{topic}' suggests further development in application Q."
    ]
    return mock_results

# --- 3. CONTENT GENERATION WITH LLM ---
def generate_slide_content_with_llm(topic, web_search_snippets, llm_provider=LLM_PROVIDER):
    """
    Generates slide content using an LLM.
    The LLM should synthesize its knowledge with the web_search_snippets.
    It MUST return a structured JSON object as defined in the prompt.
    """
    print(f"\n[LLM - Using {llm_provider}] Generating slide content for: \"{topic}\"...")

    prompt_template = f"""
    You are an expert content creator tasked with generating a structured 7-slide presentation on the topic: "{topic}".
    Incorporate your own knowledge and synthesize it with the following information from recent web search results:
    --- WEB SEARCH SNIPPETS START ---
    {json.dumps(web_search_snippets, indent=2)}
    --- WEB SEARCH SNIPPETS END ---

    The presentation structure MUST be as follows:
    - Slide 1: Title Slide (a compelling main title for the presentation)
    - Slide 2: Overview (a title for the overview and 2-3 main bullet points summarizing the presentation's scope)
    - Slide 3: Key Point 1 (a concise title for this key point and 2-4 supporting bullet points)
    - Slide 4: Key Point 2 (a concise title for this key point and 2-4 supporting bullet points, distinct from Key Point 1)
    - Slide 5: Key Point 3 (a concise title for this key point and 2-4 supporting bullet points, distinct from Key Points 1 & 2)
    - Slide 6: Key Point 4 (a concise title for this key point and 2-4 supporting bullet points, distinct from Key Points 1, 2 & 3)
    - Slide 7: Conclusion / Takeaways (a title for the conclusion and 2-3 bullet points summarizing key messages or takeaways)

    Provide the output STRICTLY as a single JSON object adhering to the following schema. Do NOT include any explanatory text before or after the JSON object.

    {{
      "slide_1_title": "string (Main Presentation Title)",
      "slide_2_overview": {{
        "title": "string (e.g., 'Overview', 'Executive Summary')",
        "points": ["string (Bullet Point 1)", "string (Bullet Point 2)"]
      }},
      "slide_3_key_point_1": {{
        "title": "string (Title for Key Point 1)",
        "points": ["string (Detail A for KP1)", "string (Detail B for KP1)"]
      }},
      "slide_4_key_point_2": {{
        "title": "string (Title for Key Point 2)",
        "points": ["string (Detail A for KP2)", "string (Detail B for KP2)"]
      }},
      "slide_5_key_point_3": {{
        "title": "string (Title for Key Point 3)",
        "points": ["string (Detail A for KP3)", "string (Detail B for KP3)"]
      }},
      "slide_6_key_point_4": {{
        "title": "string (Title for Key Point 4)",
        "points": ["string (Detail A for KP4)", "string (Detail B for KP4)"]
      }},
      "slide_7_conclusion": {{
        "title": "string (e.g., 'Conclusion', 'Key Takeaways')",
        "points": ["string (Takeaway 1)", "string (Takeaway 2)"]
      }}
    }}

    Ensure all text is concise and suitable for PowerPoint slides.
    Make sure the key points are distinct and cover different aspects of the topic.
    """

    # Using Google Gemini 
    if llm_provider == "GEMINI":
        gemini_api_key = os.getenv("GEMINI_API_KEY")
        if gemini_api_key:
            try:
                import google.generativeai as genai 
                genai.configure(api_key=gemini_api_key)
                model = genai.GenerativeModel(
                    'gemini-1.5-flash-latest', 
                    generation_config=genai.types.GenerationConfig(
                        response_mime_type="application/json" 
                    )
                )
                print(f"[LLM - {llm_provider}] Sending prompt to Gemini API...")
                response = model.generate_content(prompt_template)
                content_json_str = response.text 
                print(f"[LLM - {llm_provider}] Received raw response (first 100 chars): {content_json_str[:100]}...")
                parsed_content = json.loads(content_json_str)
                print(f"[LLM - {llm_provider}] Successfully parsed JSON content.")
                return parsed_content
            except ImportError:
                print(f"[LLM - {llm_provider} Error] The 'google-generativeai' library is not installed. Please install it using 'pip install google-generativeai'.")
            except Exception as e:
                print(f"[LLM - {llm_provider} Error] An error occurred: {e}")
                raw_response_text = "N/A"
                if 'response' in locals() and hasattr(response, 'text'):
                    raw_response_text = response.text
                elif 'response' in locals() and hasattr(response, 'parts'):
                    try:
                        raw_response_text = "".join(part.text for part in response.parts)
                    except Exception: 
                         raw_response_text = str(response.parts)
                print(f"[LLM - {llm_provider} Error] Raw response was: {raw_response_text[:500]}...")
        else:
            print(f"[LLM - {llm_provider} Error] GEMINI_API_KEY not found in .env file.")
    
    
    print(f"[LLM - Fallback] API call failed for {llm_provider} or provider not configured. Using mock slide content.")
    mock_slide_content = {
        "slide_1_title": f"A Comprehensive Analysis of {topic}",
        "slide_2_overview": {
            "title": "Presentation Overview",
            "points": [
                f"Introduction to the core concepts of {topic}.",
                f"Exploration of key areas and impacts related to {topic}.",
                "Summary of current trends and future projections."
            ]
        },
        "slide_3_key_point_1": {
            "title": f"Fundamental Aspects of {topic}",
            "points": [
                "Defining the primary characteristics and components.",
                "Historical context and evolution.",
                "Its significance in the broader field/industry."
            ]
        },
        "slide_4_key_point_2": {
            "title": f"Current Trends and Developments in {topic}",
            "points": [
                "Highlighting recent advancements and innovations.",
                f"Statistical data or notable examples of {topic} in action.",
                "Emerging patterns and shifts in understanding or application."
            ]
        },
        "slide_5_key_point_3": {
            "title": f"Challenges and Opportunities for {topic}",
            "points": [
                f"Identifying key obstacles or limitations concerning {topic}.",
                "Potential areas for growth, research, or improvement.",
                "Mitigation strategies for challenges."
            ]
        },
        "slide_6_key_point_4": {
            "title": f"The Future Outlook of {topic}",
            "points": [
                f"Predictions for the evolution of {topic} in the next 5-10 years.",
                f"Potential impact of {topic} on society, technology, or specific sectors.",
                "Upcoming research directions or anticipated breakthroughs."
            ]
        },
        "slide_7_conclusion": {
            "title": "Conclusion and Key Takeaways",
            "points": [
                f"Recap of the most critical findings about {topic}.",
                f"Final thoughts on the importance and relevance of {topic}.",
                "Recommendations for further study or action."
            ]
        }
    }
    return mock_slide_content

# --- 4. POWERPOINT SLIDE CREATION ---
def create_presentation_from_content(topic_name, content_json, template_path=None):
    """
    Creates a PowerPoint presentation using python-pptx.
    Populates slides with titles and bullet points from the content_json.
    Optionally uses a template if template_path is provided.
    """
    print("\n[PPTX] Creating PowerPoint presentation...")
    try:
        if template_path and os.path.exists(template_path):
            prs = Presentation(template_path)
            print(f"[PPTX] Using template: {template_path}")
        else:
            if template_path: 
                print(f"[PPTX] Warning: Template '{template_path}' not found. Using default layout.")
            prs = Presentation() 
    except Exception as e:
        print(f"[PPTX Error] Failed to load presentation or template: {e}. Using default layout.")
        prs = Presentation()

    def add_content_slide_with_bullets(prs, slide_layout_idx, title_text, points_list):
        """Helper function to add a content slide with a title and bullet points."""
        try:
            slide_layout = prs.slide_layouts[slide_layout_idx]
        except IndexError:
            print(f"[PPTX Warning] Slide layout index {slide_layout_idx} out of range. Using layout 1 (Title and Content).")
            slide_layout = prs.slide_layouts[1] 
            
        slide = prs.slides.add_slide(slide_layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = title_text
        else: 
            print(f"[PPTX Warning] Slide layout for '{title_text}' might not have a dedicated title placeholder. Adding text box for title.")
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
            tf = txBox.text_frame
            tf.text = title_text
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.size = Pt(28)

        body_placeholder = None
        for shape in slide.placeholders: 
            if shape.placeholder_format.idx == 1 or \
               (shape.name and ("Body" in shape.name or "Content" in shape.name or "Object" in shape.name)): 
                if shape.has_text_frame:
                    body_placeholder = shape
                    break
        if not body_placeholder and slide.placeholders: 
             for shape in slide.placeholders:
                 if shape.has_text_frame and getattr(slide.shapes, 'title', None) != shape:
                     body_placeholder = shape
                     break

        if body_placeholder:
            tf = body_placeholder.text_frame
            tf.clear() 
            tf.word_wrap = True
            for point_text in points_list:
                p = tf.add_paragraph()
                p.text = str(point_text)
                p.level = 0 
                p.font.size = Pt(18)
        else: 
            print(f"[PPTX Warning] Could not find a suitable body placeholder for slide '{title_text}'. Adding a new textbox for bullets.")
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(5.0))
            tf = txBox.text_frame
            tf.word_wrap = True
            for point_text in points_list:
                p = tf.add_paragraph()
                p.text = str(point_text)
                p.level = 0
                p.font.size = Pt(18)
        print(f"[PPTX] Added Slide: {title_text}")

    title_slide_layout_idx = 0 
    content_slide_layout_idx = 1

    slide1_title_text = content_json.get("slide_1_title", f"{topic_name} - An Overview")
    try:
        title_slide_layout = prs.slide_layouts[title_slide_layout_idx]
        slide = prs.slides.add_slide(title_slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = slide1_title_text
        
        subtitle_placeholder = None 
        if len(slide.placeholders) > 1:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1 or (ph.name and "Subtitle" in ph.name):
                    if ph.has_text_frame: subtitle_placeholder = ph; break
            if not subtitle_placeholder and slide.placeholders[1].has_text_frame: 
                 subtitle_placeholder = slide.placeholders[1]

        if subtitle_placeholder:
            subtitle_placeholder.text = f"AI-Generated Presentation on: {topic_name}"
            if subtitle_placeholder.text_frame.paragraphs:
                 subtitle_placeholder.text_frame.paragraphs[0].font.size = Pt(18)
            else:
                p = subtitle_placeholder.text_frame.add_paragraph()
                p.text = f"AI-Generated Presentation on: {topic_name}"; p.font.size = Pt(18)
        print(f"[PPTX] Added Slide 1: {slide1_title_text}")
    except Exception as e: 
        print(f"[PPTX Error] Failed to create title slide using layout {title_slide_layout_idx}: {e}. Attempting fallback.")
        add_content_slide_with_bullets(prs, content_slide_layout_idx, slide1_title_text, [f"AI-Generated Presentation on: {topic_name}"])

    overview_data = content_json.get("slide_2_overview", {})
    add_content_slide_with_bullets(prs, content_slide_layout_idx,
                                   overview_data.get("title", "Overview"),
                                   overview_data.get("points", ["No overview points generated."]))

    for i in range(1, 5): 
        key_point_data = content_json.get(f"slide_{i+2}_key_point_{i}", {}) 
        add_content_slide_with_bullets(prs, content_slide_layout_idx,
                                       key_point_data.get("title", f"Key Point {i}"),
                                       key_point_data.get("points", [f"No points generated for Key Point {i}."]))
        
    conclusion_data = content_json.get("slide_7_conclusion", {})
    add_content_slide_with_bullets(prs, content_slide_layout_idx,
                                   conclusion_data.get("title", "Conclusion / Takeaways"),
                                   conclusion_data.get("points", ["No conclusion points generated."]))

    clean_file_name_base = re.sub(r'[^\w\-.]', '_', topic_name.lower()) 
    if not clean_file_name_base: clean_file_name_base = "presentation"
    file_name = f"{clean_file_name_base}_presentation.pptx"

    try:
        prs.save(file_name)
        print(f"\n[PPTX] Presentation saved successfully as: {file_name}")
        return file_name
    except Exception as e:
        print(f"[PPTX Error] Failed to save presentation '{file_name}': {e}")
        fallback_file_name = "fallback_presentation.pptx"
        try:
            prs.save(fallback_file_name)
            print(f"[PPTX] Presentation saved with fallback name: {fallback_file_name}")
            return fallback_file_name
        except Exception as e2:
            print(f"[PPTX Error] Failed to save with fallback name either: {e2}")
            return None

# --- MAIN EXECUTION ---
if __name__ == "__main__":
    print("--- Automated Slide Deck Generator ---")
    
    topic = get_topic_from_user()
    web_snippets = perform_web_search(topic) 
    if not web_snippets:
        print("[Main] Warning: Web search returned no snippets. LLM will rely on its own knowledge.")
        web_snippets = [f"No specific web information found for {topic}, relying on general knowledge."]

    slide_content_json = generate_slide_content_with_llm(topic, web_snippets, llm_provider=LLM_PROVIDER)

    if not slide_content_json:
        print("\n[Main Error] Failed to generate slide content from LLM. Exiting.")
    else:
        custom_template_path = None 
        
        presentation_file = create_presentation_from_content(topic, slide_content_json, template_path=custom_template_path)
        
        if presentation_file:
            print(f"\nSuccessfully generated presentation: {os.path.abspath(presentation_file)}")
        else:
            print("\n[Main Error] Failed to create or save the presentation file.")

    print("\n--- Script Finished ---")